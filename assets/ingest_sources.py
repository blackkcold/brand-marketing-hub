#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Normalize DOCX/XLSX/PPTX/PDF/CSV into v4.1 source + content-unit manifests.

Host-native Files/Documents/Spreadsheets/Presentations capabilities are preferred.
This CLI fallback preserves text, tables, spreadsheet formulas, charts/notes where
accessible, and embedded source images without leaking absolute local file paths.
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import mimetypes
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

SUPPORTED={".docx":"docx",".xlsx":"xlsx",".pptx":"pptx",".pdf":"pdf",".csv":"csv"}

def sha256(path:Path)->str:
    h=hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda:f.read(1024*1024),b""):
            h.update(chunk)
    return h.hexdigest()

def sha256_blob(blob:bytes)->str:
    return hashlib.sha256(blob).hexdigest()

class UnitBuilder:
    def __init__(self,out_dir:Path):
        self.n=0
        self.units=[]
        self.out_dir=out_dir
        self.asset_dir=out_dir/"extracted_assets"
        self.asset_dir.mkdir(parents=True,exist_ok=True)
        self.image_n=0

    def add(self,source_id:str,unit_type:str,locator:dict[str,Any],content:Any,
            preserve_mode="summarize-ok",importance="normal",formulae=None,notes=None):
        self.n+=1
        item={
            "unit_id":f"UNIT-{self.n:04d}",
            "source_id":source_id,
            "unit_type":unit_type,
            "locator":locator,
            "content":content,
            "preserve_mode":preserve_mode,
            "importance":importance,
        }
        if formulae: item["formulae"]=formulae
        if notes: item["notes"]=notes
        self.units.append(item)

    def add_image(self,source_id:str,locator:dict[str,Any],blob:bytes,ext:str,
                  content_type=None,importance="normal",notes=None):
        if not blob:return
        self.image_n+=1
        ext=(ext or "bin").lstrip(".").lower()
        out=self.asset_dir/f"{source_id.lower()}-img-{self.image_n:04d}.{ext}"
        out.write_bytes(blob)
        rel=out.relative_to(self.out_dir).as_posix()
        content={
            "extracted_path":rel,
            "sha256":sha256_blob(blob),
            "content_type":content_type,
            "bytes":len(blob),
        }
        self.add(source_id,"image",locator,content,preserve_mode="semantic",
                 importance=importance,notes=notes)

def ingest_docx(path:Path,source_id:str,b:UnitBuilder):
    doc=Document(path)
    para_no=0
    for p in doc.paragraphs:
        text=p.text.strip()
        if not text:continue
        para_no+=1
        style=(p.style.name or "").lower() if p.style else ""
        is_heading=style.startswith("heading") or style.startswith("标题")
        b.add(source_id,"heading" if is_heading else "paragraph",
              {"paragraph":para_no},text,
              preserve_mode="semantic" if is_heading else "summarize-ok",
              importance="high" if is_heading else "normal")
    for ti,table in enumerate(doc.tables,1):
        rows=[[cell.text for cell in row.cells] for row in table.rows]
        b.add(source_id,"table",{"table":ti},{"rows":rows},
              preserve_mode="semantic",importance="high")

    image_no=0
    for rel in doc.part.rels.values():
        if not str(rel.reltype).endswith("/image"):continue
        try:
            part=rel.target_part
            blob=part.blob
            ext=Path(str(part.partname)).suffix.lstrip(".") or "bin"
            content_type=getattr(part,"content_type",None)
            image_no+=1
            b.add_image(source_id,{"image":image_no},blob,ext,content_type,
                        notes="Embedded DOCX image; paragraph-level anchor is not exposed by fallback parser.")
        except Exception:
            continue

def _xlsx_value(cell):
    v=cell.value
    if v is None:return None
    if hasattr(v,"isoformat"):
        try:return v.isoformat()
        except Exception:pass
    return v

def _xlsx_anchor(image):
    try:
        marker=image.anchor._from
        return f"{marker.col+1},{marker.row+1}"
    except Exception:
        return None

def ingest_xlsx(path:Path,source_id:str,b:UnitBuilder):
    wb_formula=load_workbook(path,data_only=False,read_only=False)
    wb_values=load_workbook(path,data_only=True,read_only=False)
    for ws in wb_formula.worksheets:
        if ws.sheet_state!="visible":continue
        min_r=min_c=None; max_r=max_c=None
        formulas=[]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    min_r=cell.row if min_r is None else min(min_r,cell.row)
                    max_r=cell.row if max_r is None else max(max_r,cell.row)
                    min_c=cell.column if min_c is None else min(min_c,cell.column)
                    max_c=cell.column if max_c is None else max(max_c,cell.column)
        if min_r is not None:
            from openpyxl.utils import get_column_letter
            range_ref=f"{get_column_letter(min_c)}{min_r}:{get_column_letter(max_c)}{max_r}"
            rows=[]
            ws_values=wb_values[ws.title]
            for r in range(min_r,max_r+1):
                row=[]
                for col in range(min_c,max_c+1):
                    fcell=ws.cell(r,col); vcell=ws_values.cell(r,col)
                    row.append(_xlsx_value(vcell))
                    if fcell.data_type=="f":
                        formulas.append({"cell":fcell.coordinate,"formula":str(fcell.value),
                                         "cached_value":_xlsx_value(vcell)})
                rows.append(row)
            b.add(source_id,"cell-range",{"sheet":ws.title,"range":range_ref},
                  {"rows":rows,"dimensions":{"rows":len(rows),"cols":max_c-min_c+1}},
                  preserve_mode="semantic",importance="high",formulae=formulas or None)

        for ii,img in enumerate(getattr(ws,"_images",[]),1):
            try:
                blob=img._data()
                ext=(getattr(img,"format",None) or "png").lower()
                anchor=_xlsx_anchor(img)
                locator={"sheet":ws.title,"image":ii}
                if anchor:locator["range"]=anchor
                b.add_image(source_id,locator,blob,ext,
                            importance="normal",notes="Embedded XLSX image.")
            except Exception:
                continue

        for ci,chart in enumerate(getattr(ws,"_charts",[]),1):
            title=None
            try:
                title=str(chart.title.tx.rich.p[0].r[0].t)
            except Exception:
                pass
            b.add(source_id,"metadata",{"sheet":ws.title,"shape":ci},
                  {"object":"chart","chart_type":chart.__class__.__name__,"title":title},
                  preserve_mode="semantic",importance="normal",
                  notes="Workbook chart metadata; underlying values remain preserved in cell-range units.")

def ingest_pptx(path:Path,source_id:str,b:UnitBuilder):
    prs=Presentation(path)
    master_meta=[]
    for mi,master in enumerate(prs.slide_masters,1):
        master_meta.append({
            "master":mi,
            "layout_names":[layout.name for layout in master.slide_layouts],
            "shape_count":len(master.shapes),
        })
    b.add(source_id,"metadata",{},
          {"object":"presentation","slide_count":len(prs.slides),
           "slide_width":prs.slide_width,"slide_height":prs.slide_height,
           "masters":master_meta},
          preserve_mode="optional",importance="low",
          notes="PPTX presentation/master inventory; native runtime should inspect full theme/master when this source is a style-reference.")
    for si,slide in enumerate(prs.slides,1):
        texts=[]
        for sh in slide.shapes:
            if getattr(sh,"has_text_frame",False):
                txt=sh.text_frame.text.strip()
                if txt:texts.append({"shape":sh.shape_id,"text":txt})
            if getattr(sh,"has_table",False):
                rows=[[cell.text for cell in row.cells] for row in sh.table.rows]
                b.add(source_id,"slide-table",{"slide":si,"shape":sh.shape_id},
                      {"rows":rows},preserve_mode="semantic",importance="high")
            if getattr(sh,"has_chart",False):
                chart=sh.chart; series=[]
                try:
                    for s in chart.series:
                        vals=list(s.values) if hasattr(s,"values") else []
                        series.append({"name":getattr(s,"name",None),"values":vals})
                except Exception:
                    pass
                b.add(source_id,"slide-chart",{"slide":si,"shape":sh.shape_id},
                      {"chart_type":str(chart.chart_type),"series":series},
                      preserve_mode="semantic",importance="high")
            if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
                try:
                    image=sh.image
                    b.add_image(source_id,{"slide":si,"shape":sh.shape_id},image.blob,image.ext,
                                image.content_type,importance="normal",
                                notes="Embedded PPTX picture.")
                except Exception:
                    pass
        if texts:
            b.add(source_id,"slide-text",{"slide":si},{"shapes":texts},
                  preserve_mode="semantic",importance="high")
        try:
            notes=slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                b.add(source_id,"speaker-notes",{"slide":si},notes,
                      preserve_mode="summarize-ok",importance="normal")
        except Exception:
            pass

def ingest_pdf(path:Path,source_id:str,b:UnitBuilder):
    reader=PdfReader(path)
    for pi,page in enumerate(reader.pages,1):
        text=(page.extract_text() or "").strip()
        b.add(source_id,"pdf-page",{"page":pi},text,
              preserve_mode="summarize-ok",importance="normal",
              notes=None if text else "No extractable text; host visual page inspection is required.")

def ingest_csv(path:Path,source_id:str,b:UnitBuilder):
    with path.open("r",encoding="utf-8-sig",newline="") as f:
        rows=list(csv.reader(f))
    b.add(source_id,"cell-range",{"sheet":"CSV","range":f"A1:{len(rows)}"},
          {"rows":rows},preserve_mode="semantic",importance="high")

INGESTORS={"docx":ingest_docx,"xlsx":ingest_xlsx,"pptx":ingest_pptx,"pdf":ingest_pdf,"csv":ingest_csv}

def main()->int:
    ap=argparse.ArgumentParser()
    ap.add_argument("inputs",nargs="+",type=Path)
    ap.add_argument("--out-dir",type=Path,default=Path("source_bundle"))
    ap.add_argument("--confidentiality",choices=["public","internal","confidential","restricted"],default="internal")
    ap.add_argument("--role",action="append",default=[],metavar="FILE=ROLE",
                    help="Override a source role: content-source|data-source|style-reference|evidence-source|mixed")
    args=ap.parse_args()
    valid_roles={"content-source","data-source","style-reference","evidence-source","mixed"}
    role_overrides={}
    for item in args.role:
        if "=" not in item:
            raise SystemExit(f"invalid --role value: {item}; expected FILE=ROLE")
        name,role=item.split("=",1)
        if role not in valid_roles:
            raise SystemExit(f"invalid source role: {role}")
        role_overrides[Path(name).name]=role
    args.out_dir.mkdir(parents=True,exist_ok=True)
    b=UnitBuilder(args.out_dir)
    sources=[]
    for i,path in enumerate(args.inputs,1):
        if not path.exists():raise SystemExit(f"missing input: {path}")
        kind=SUPPORTED.get(path.suffix.lower())
        if not kind:raise SystemExit(f"unsupported input type: {path.suffix} ({path})")
        sid=f"FILE-{i:03d}"
        sources.append({
            "source_id":sid,
            "kind":kind,
            "display_name":path.name,
            "origin":"user-upload",
            "mime_type":mimetypes.guess_type(path.name)[0],
            "path":path.name,
            "url":None,
            "sha256":sha256(path),
            "confidentiality":args.confidentiality,
            "role":role_overrides.get(path.name, "data-source" if kind in ("xlsx","csv") else "content-source")
        })
        INGESTORS[kind](path,sid,b)
    (args.out_dir/"source_inventory.json").write_text(
        json.dumps({"version":"4.1","sources":sources},ensure_ascii=False,indent=2),encoding="utf-8")
    (args.out_dir/"content_units.json").write_text(
        json.dumps({"version":"4.1","units":b.units},ensure_ascii=False,indent=2),encoding="utf-8")
    print(f"PASS: {len(sources)} sources -> {len(b.units)} content units in {args.out_dir}")
    return 0

if __name__=="__main__":
    raise SystemExit(main())
