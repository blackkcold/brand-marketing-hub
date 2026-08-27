#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""v4.1 PPTX structural/brand validator.

Checks delivery-blocking PPTX defects that can be verified deterministically:
placeholders, out-of-bounds shapes, external linked images, recognizable vivo
wordmark presence, template font/palette contract and basic legacy structure.
Visual/semantic fidelity still requires host render + visual review.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PLACEHOLDER_RE=re.compile(r'【[^】]*】')
ROOT=Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE=ROOT/"brand/vivo/template-manifest.json"
KNOWN_LOGO=ROOT/"assets/vivo-deck/vivo_wordmark_white.png"

SEVERITY={
 "page_count_mismatch":"P0",
 "page_numbers_not_contiguous":"P0",
 "placeholder_remaining":"P0",
 "external_image_relationship":"P0",
 "table_row_count_exceeded":"P0",
 "table_col_count_exceeded":"P0",
 "source_missing":"P1",
 "shape_out_of_bounds":"P1",
 "logo_missing":"P1",
 "bullet_count_exceeded":"P1",
 "card_count_exceeded":"P1",
 "font_off_contract":"P2",
 "color_off_palette":"P2",
 "partner_color_requires_verification":"P3",
 "shadow_forbidden":"P2",
 "page_num_missing":"P2",
 "sub_too_long":"P2",
 "text_overlap":"P1",
 "font_too_small":"P1",
 "canvas_not_16_9":"P1",
}

def issue(code,**kwargs):
    return {"severity":SEVERITY.get(code,"P2"),"code":code,**kwargs}

def sha256_bytes(data:bytes)->str:
    return hashlib.sha256(data).hexdigest()

def file_sha256(path:Path)->str|None:
    if not path.exists(): return None
    return sha256_bytes(path.read_bytes())

def shape_text(shape):
    if getattr(shape,"has_text_frame",False): return shape.text_frame.text
    if getattr(shape,"has_table",False):
        return "\n".join(cell.text for row in shape.table.rows for cell in row.cells)
    return ""

def md_pages(path):
    pages=[]; current=None
    for line in Path(path).read_text(encoding="utf-8").splitlines():
        m=re.match(r'^##\s+P(\d+)[｜|]\s*(.+)$',line)
        if m:
            current={"no":int(m.group(1)),"title":m.group(2).strip(),"layout":None,"sub":None,"source":None,
                     "bullets":0,"cards":0,"table":0,"table_cols":0}
            pages.append(current); continue
        if current is None: continue
        if line.startswith("@layout"):
            parts=line.split(); current["layout"]=parts[1] if len(parts)>1 else None
        elif line.startswith("@sub"): current["sub"]=line[4:].strip()
        elif line.startswith("@source"): current["source"]=line[7:].strip()
        elif line.startswith("### "): current["cards"]+=1
        elif re.match(r'^-\s+',line): current["bullets"]+=1
        elif line.strip().startswith("|") and not set(line.strip().strip("|").replace("|","").strip())<=set("-:"):
            current["table"]+=1
            current["table_cols"]=max(current["table_cols"],len(line.strip().strip("|").split("|")))
    return pages

def load_template(path:Path|None):
    p=path or DEFAULT_TEMPLATE
    if not p.exists(): return {"allowed_fonts":[],"allowed_palettes":{},"partner_accent_policy":"forbid"}
    return json.loads(p.read_text(encoding="utf-8"))

def palette_set(template):
    out=set()
    for values in template.get("allowed_palettes",{}).values():
        if isinstance(values,list): out.update(str(v).upper().replace("#","") for v in values)
    return out

def rgb_hex(color):
    try:return str(color).upper()
    except Exception:return None

def external_image_relationships(pptx_path):
    found=[]
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if not name.endswith(".rels"): continue
            try:root=ET.fromstring(z.read(name))
            except Exception:continue
            for rel in root:
                typ=rel.attrib.get("Type","")
                mode=rel.attrib.get("TargetMode","")
                if typ.endswith("/image") and mode.lower()=="external":
                    found.append({"rels":name,"target":rel.attrib.get("Target")})
    return found

def picture_hash(shape):
    if shape.shape_type!=MSO_SHAPE_TYPE.PICTURE:return None
    try:return sha256_bytes(shape.image.blob)
    except Exception:return None

def has_known_logo(prs):
    known=file_sha256(KNOWN_LOGO)
    if not known:return False
    for slide in prs.slides:
        for sh in slide.shapes:
            if picture_hash(sh)==known:return True
    try:
        for master in prs.slide_masters:
            for sh in master.shapes:
                if picture_hash(sh)==known:return True
    except Exception:pass
    return False

def bbox_in(shape):
    try:
        l,t,w,h=[v/914400 for v in (shape.left,shape.top,shape.width,shape.height)]
        return (l,t,l+w,t+h)
    except Exception:
        return None

def overlap_ratio(a,b):
    ax1,ay1,ax2,ay2=a; bx1,by1,bx2,by2=b
    iw=max(0,min(ax2,bx2)-max(ax1,bx1))
    ih=max(0,min(ay2,by2)-max(ay1,by1))
    inter=iw*ih
    aa=max(0,ax2-ax1)*max(0,ay2-ay1)
    ba=max(0,bx2-bx1)*max(0,by2-by1)
    denom=min(aa,ba)
    return inter/denom if denom>0 else 0

def validate(pptx_path,md_path=None,template_manifest=None):
    prs=Presentation(pptx_path)
    template=load_template(Path(template_manifest) if template_manifest else None)
    allowed_fonts=set(template.get("allowed_fonts",[]))
    allowed_colors=palette_set(template)
    partner_policy=template.get("partner_accent_policy","forbid")
    issues=[]
    slide_w=prs.slide_width/914400
    slide_h=prs.slide_height/914400
    ratio=slide_w/slide_h if slide_h else 0
    if abs(ratio-(16/9))>0.02:
        issues.append(issue("canvas_not_16_9",width=round(slide_w,3),height=round(slide_h,3),ratio=round(ratio,4)))
    info={"version":"4.1","file":pptx_path,"slide_count":len(prs.slides),
          "canvas_in":[round(slide_w,3),round(slide_h,3)],
          "issues":[],"template_id":template.get("template_id")}

    for rel in external_image_relationships(pptx_path):
        issues.append(issue("external_image_relationship",target=rel["target"],rels=rel["rels"]))

    source_pages=md_pages(md_path) if md_path and os.path.exists(md_path) else []
    if source_pages:
        if len(prs.slides)!=len(source_pages):
            issues.append(issue("page_count_mismatch",expected=len(source_pages),actual=len(prs.slides)))
        nums=[p["no"] for p in source_pages]
        if nums!=list(range(1,len(nums)+1)):
            issues.append(issue("page_numbers_not_contiguous",numbers=nums))
        for p in source_pages:
            page=p["no"]
            if p["sub"] and len(p["sub"])>22: issues.append(issue("sub_too_long",page=page,length=len(p["sub"])))
            if p["bullets"]>6 and p["layout"] in (None,"bullets","purpose"):
                issues.append(issue("bullet_count_exceeded",page=page,count=p["bullets"]))
            if p["cards"]>8: issues.append(issue("card_count_exceeded",page=page,count=p["cards"]))
            if p["layout"] in ("table","budget","matrix","timeline") and p["table"]>10:
                issues.append(issue("table_row_count_exceeded",page=page,count=p["table"]))
            if p["layout"] in ("table","budget","matrix","timeline") and p["table_cols"]>6:
                issues.append(issue("table_col_count_exceeded",page=page,count=p["table_cols"]))
            if p["layout"] in ("stats","evidence-grid") and not p["source"]:
                issues.append(issue("source_missing",page=page))

    logo_present=has_known_logo(prs)
    if not logo_present and len(prs.slide_masters)<2:
        issues.append(issue("logo_missing",message="No known embedded vivo wordmark found; master/reference branding must be visually verified."))

    for idx,slide in enumerate(prs.slides):
        slide_no=idx+1; cover_or_end=idx in (0,len(prs.slides)-1); has_page=False
        text_boxes=[]
        for shape in slide.shapes:
            box=bbox_in(shape)
            text=shape_text(shape)
            material=bool(text.strip()) or getattr(shape,"has_table",False) or getattr(shape,"has_chart",False)
            if box and material:
                l,t,r,b=box
                if l<-.01 or t<-.01 or r>slide_w+.01 or b>slide_h+.01:
                    issues.append(issue("shape_out_of_bounds",slide=slide_no,shape=shape.shape_id))
            if box and getattr(shape,"has_text_frame",False) and text.strip():
                text_boxes.append((shape,box,text.strip()))
            if getattr(shape,"has_text_frame",False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and allowed_fonts and run.font.name not in allowed_fonts:
                            issues.append(issue("font_off_contract",slide=slide_no,shape=shape.shape_id,font=run.font.name))
                        if run.text.strip() and run.font.size and run.font.size.pt < 7:
                            issues.append(issue("font_too_small",slide=slide_no,shape=shape.shape_id,
                                                font_size=round(run.font.size.pt,2),text=run.text[:40]))
            try:
                if "outerShdw" in shape._element.xml or "innerShdw" in shape._element.xml:
                    issues.append(issue("shadow_forbidden",slide=slide_no,shape=shape.shape_id))
            except Exception:pass
            if re.fullmatch(r'\d+',text.strip()): has_page=True
            for m in PLACEHOLDER_RE.finditer(text):
                issues.append(issue("placeholder_remaining",slide=slide_no,token=m.group(0)))
            if getattr(shape,"fill",None) and shape.fill.type:
                try:
                    value=rgb_hex(shape.fill.fore_color.rgb)
                    if value and allowed_colors and value not in allowed_colors and value!="00000000":
                        code="partner_color_requires_verification" if partner_policy=="allow-verified-brand-colors" else "color_off_palette"
                        issues.append(issue(code,slide=slide_no,shape=shape.shape_id,color=value))
                except Exception:pass
        for i in range(len(text_boxes)):
            shape_a,box_a,text_a=text_boxes[i]
            if len(text_a)<=1 or re.fullmatch(r'\d+',text_a): continue
            for j in range(i+1,len(text_boxes)):
                shape_b,box_b,text_b=text_boxes[j]
                if len(text_b)<=1 or re.fullmatch(r'\d+',text_b): continue
                if overlap_ratio(box_a,box_b)>=0.35:
                    issues.append(issue("text_overlap",slide=slide_no,
                                        shapes=[shape_a.shape_id,shape_b.shape_id],
                                        texts=[text_a[:60],text_b[:60]]))
        if not cover_or_end and not has_page:
            issues.append(issue("page_num_missing",slide=slide_no))

    blocking=[x for x in issues if x["severity"] in ("P0","P1")]
    info["issues"]=issues
    info["counts_by_severity"]={lvl:sum(1 for x in issues if x["severity"]==lvl) for lvl in ("P0","P1","P2","P3")}
    info["issue_policy"]="v4.1 severity gate: any P0/P1 blocks delivery"
    info["pass"]=not issues
    info["output_allowed"]=not blocking
    info["status"]="pass" if not issues else ("fail" if blocking else "pass-with-cosmetic-notes")
    return info

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument("pptx"); ap.add_argument("--md"); ap.add_argument("--template-manifest"); ap.add_argument("--json",action="store_true")
    args=ap.parse_args()
    if not os.path.exists(args.pptx):
        print(json.dumps({"error":"pptx 不存在","path":args.pptx},ensure_ascii=False)); return 2
    info=validate(args.pptx,args.md,args.template_manifest)
    if args.json: print(json.dumps(info,ensure_ascii=False,indent=2))
    else:
        print(f'文件: {info["file"]}\n页数: {info["slide_count"]}\n状态: {info["status"]}')
        for item in info["issues"]: print("  - "+json.dumps(item,ensure_ascii=False))
    return 0 if info["output_allowed"] else 1

if __name__=="__main__":
    raise SystemExit(main())
