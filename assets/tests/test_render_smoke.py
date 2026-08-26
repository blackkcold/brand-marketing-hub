#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Non-destructive smoke test for the brand-marketing-hub PPT renderer.

The harness intentionally keeps all generated inputs/outputs in a TemporaryDirectory.
It coordinates with the existing renderer CLI without editing renderer code.
"""

from __future__ import annotations

import re
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SKILL_ROOT = Path(__file__).resolve().parents[2]
ASSETS_DIR = SKILL_ROOT / "assets"
FIXTURE = ASSETS_DIR / "fixtures" / "visual-regression.md"
RENDERER = ASSETS_DIR / "md2pptx_vivo.py"
EXPECTED_SLIDES = 15
EMU_PER_INCH = 914400


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def make_synthetic_image(path: Path) -> None:
    img = Image.new("RGB", (960, 540), "#1E46E6")
    draw = ImageDraw.Draw(img)
    for x in range(0, 960, 24):
        shade = 70 + (x * 120 // 960)
        draw.rectangle([x, 0, x + 24, 540], fill=(30, shade, 230))
    draw.rectangle([96, 96, 864, 444], outline="white", width=6)
    draw.text((130, 245), "synthetic visual fixture", fill="white")
    img.save(path)


def materialize_fixture(tmpdir: Path, *, sanitize_placeholders: bool) -> Path:
    image_path = tmpdir / "synthetic-fixture-image.png"
    make_synthetic_image(image_path)

    text = read_text(FIXTURE).replace("__VISUAL_FIXTURE_IMAGE__", str(image_path))
    if sanitize_placeholders:
        text = re.sub(r"【[^】]+】", "已替换合成字段", text)

    rendered_input = tmpdir / "visual-regression.renderable.md"
    rendered_input.write_text(text, encoding="utf-8")
    return rendered_input


def parse_pages(md_text: str) -> list[dict[str, object]]:
    pages: list[dict[str, object]] = []
    current: dict[str, object] | None = None
    for raw in md_text.splitlines():
        line = raw.rstrip("\n")
        match = re.match(r"^##\s+P(\d+)[｜|]\s*(.+)$", line)
        if match:
            current = {
                "no": int(match.group(1)),
                "title": match.group(2).strip(),
                "layout": None,
                "layout_arg": None,
                "cards": 0,
            }
            pages.append(current)
            continue
        if current is None:
            continue
        if line.startswith("@layout"):
            parts = line.split()
            current["layout"] = parts[1] if len(parts) > 1 else None
            current["layout_arg"] = parts[2] if len(parts) > 2 else None
        elif line.startswith("### "):
            current["cards"] = int(current["cards"]) + 1
    return pages


def validate_markdown(md_path: Path, *, pptx_path: Path | None = None) -> tuple[bool, list[str]]:
    text = read_text(md_path)
    pages = parse_pages(text)
    errors: list[str] = []

    numbers = [int(p["no"]) for p in pages]
    if numbers != list(range(1, EXPECTED_SLIDES + 1)):
        errors.append(f"unexpected page numbers: {numbers}")

    cover_pages = [p for p in pages if str(p["title"]).strip() == "封面"]
    if len(cover_pages) != 1 or int(cover_pages[0]["no"]) != 1:
        errors.append("fixture must contain exactly one explicit P1 cover")

    if re.search(r"【[^】]+】", text):
        errors.append("unresolved placeholder token found")

    split_args = [p["layout_arg"] for p in pages if p["layout"] == "split"]
    if "left" not in split_args or "right" not in split_args:
        errors.append(f"split directions missing: {split_args}")

    card_counts = {int(p["no"]): int(p["cards"]) for p in pages if p["layout"] == "cards"}
    if card_counts.get(5) != 3 or card_counts.get(6) != 4:
        errors.append(f"unexpected card counts: {card_counts}")

    if pptx_path is not None:
        prs = Presentation(pptx_path)
        if len(prs.slides) != EXPECTED_SLIDES:
            errors.append(f"unexpected slide count: {len(prs.slides)}")
        # Duplicate cover regression guard: a P1 cover fixture should render to the
        # same count as markdown pages, not markdown pages + auto cover.
        if len(prs.slides) == EXPECTED_SLIDES + 1:
            errors.append("duplicate cover slide detected")

        for slide_no in (7, 8):
            slide = prs.slides[slide_no - 1]
            if not any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
                errors.append(f"slide {slide_no} split fixture image missing")
            text_widths = [
                shape.width / EMU_PER_INCH
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and "•" in shape.text_frame.text
            ]
            if not text_widths or max(text_widths) > 7.0:
                errors.append(f"slide {slide_no} split text column not constrained: {text_widths}")

        missing_slide = prs.slides[12 - 1]
        missing_text = "\n".join(
            shape.text_frame.text
            for shape in missing_slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        if "MISSING IMAGE" not in missing_text:
            errors.append("missing-image split placeholder slot not rendered")

    return not errors, errors


def run_renderer(md_path: Path, pptx_path: Path) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(RENDERER), str(md_path), str(pptx_path)],
        cwd=str(SKILL_ROOT),
        check=False,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )


def main() -> int:
    assert FIXTURE.exists(), f"missing fixture: {FIXTURE}"
    assert RENDERER.exists(), f"missing renderer: {RENDERER}"

    with tempfile.TemporaryDirectory(prefix="bmh-visual-smoke-") as tmp:
        tmpdir = Path(tmp)

        failing_md = materialize_fixture(tmpdir, sanitize_placeholders=False)
        ok, errors = validate_markdown(failing_md)
        assert not ok, "validator should fail before placeholder sanitization"
        assert any("placeholder" in error for error in errors), errors

        passing_md = materialize_fixture(tmpdir, sanitize_placeholders=True)
        output_pptx = tmpdir / "visual-regression-output.pptx"
        result = run_renderer(passing_md, output_pptx)
        assert result.returncode == 0, result.stdout + result.stderr
        assert output_pptx.exists(), "renderer did not create pptx output"

        ok, errors = validate_markdown(passing_md, pptx_path=output_pptx)
        assert ok, "; ".join(errors)

    print("PASS visual smoke: fixture renders and structural validator checks pass/fail behavior")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
