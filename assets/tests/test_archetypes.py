#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Smoke coverage for the v3 content-led page archetypes."""
from pathlib import Path
import subprocess, sys, tempfile
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
RENDERER = ROOT / 'assets' / 'md2pptx_vivo.py'

MD = '''# Archetype QA
<!-- deck: v3 layout coverage -->
<!-- deck: meta: QA｜2026｜Internal -->
## P1｜封面
## P2｜目录
- 01 证据
- 02 策略
## P3｜数据证明机会存在
@layout stats 3
@source synthetic QA data
@stat 2亿+ | 触达
@stat 68% | 重合
@stat 4.2 | 适配
## P4｜策略由四个判断组成
@layout framework
### RTB
- 可验证能力
### WHAT
- 内容动作
### WHY
- 人群需求
### HOW
- 执行机制
## P5｜合作链路需要三次推进
@layout funnel
### A
- 认知触点
### A-I
- 互动触点
### I-P
- 购买触点
## P6｜资源需要按优先级管理
@layout matrix
| 人选 | 适配 | 风险 |
|---|---|---|
| A | 高 | 低 |
## P7｜节点决定执行顺序
@layout timeline
@source synthetic QA schedule
| 阶段 | 时间 | 动作 |
|---|---|---|
| 预热 | T-7 | 预告 |
## P8｜趋势需要原生图表表达
@layout chart
@chart line
@source synthetic QA chart data
| 月份 | 指标 |
|---|---:|
| 1月 | 10 |
| 2月 | 20 |
## P9｜下一步确认决策
@layout end
'''

def main():
    with tempfile.TemporaryDirectory(prefix='bmh-archetypes-') as d:
        d = Path(d); md = d / 'input.md'; out = d / 'output.pptx'; md.write_text(MD, encoding='utf-8')
        result = subprocess.run([sys.executable, str(RENDERER), str(md), str(out)], text=True, capture_output=True)
        assert result.returncode == 0, result.stdout + result.stderr
        prs = Presentation(out)
        assert len(prs.slides) == 9
        assert any(shape.has_chart for shape in prs.slides[7].shapes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            assert run.font.name == '微软雅黑', (slide.slide_id, run.text, run.font.name)
                assert 'outerShdw' not in shape._element.xml
        end_text = '\n'.join(shape.text for shape in prs.slides[-1].shapes if getattr(shape, 'has_text_frame', False))
        assert 'THANK YOU' in end_text and '谢谢' in end_text
    print('PASS archetypes: v3 layouts render successfully')

if __name__ == '__main__': raise SystemExit(main())
