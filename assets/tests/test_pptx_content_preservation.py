#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Assert that deterministic fallback rendering does not silently drop typed source content."""
from __future__ import annotations
import argparse, json
from pathlib import Path
from pptx import Presentation

def expected_strings(spec):
    out=[]
    for slide in spec.get("slides",[]):
        for b in slide.get("content",[]):
            t=b.get("type")
            if t in ("headline","callout"):
                if b.get("text"): out.append(str(b["text"]))
            elif t=="body_text":
                value=str(b.get("text",""))
                if len(value)>500:
                    out.extend(value.split())
                elif value:
                    out.append(value)
            elif t=="bullets":
                out.extend(str(x) for x in b.get("items",[]))
            elif t=="stat":
                out.append(str(b.get("value",""))+str(b.get("unit","")))
                out.append(str(b.get("label","")))
            elif t=="table":
                for x in b.get("headers",[]): out.append(str(x))
                for row in b.get("rows",[]):
                    out.extend(str(x) for x in row)
            elif t in ("comparison","timeline"):
                for item in b.get("items",[]):
                    for v in item.values():
                        if isinstance(v,(str,int,float)): out.append(str(v))
            elif t=="source_footer":
                out.extend(str(x) for x in b.get("source_ids",[]))
    return [x for x in out if x]

def ppt_text(path):
    prs=Presentation(path)
    chunks=[]
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape,"has_text_frame",False):
                chunks.append(shape.text_frame.text)
            if getattr(shape,"has_table",False):
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
    return "\n".join(chunks),len(prs.slides)

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument("deck_spec",type=Path)
    ap.add_argument("pptx",type=Path)
    args=ap.parse_args()
    spec=json.loads(args.deck_spec.read_text(encoding="utf-8"))
    text,slide_count=ppt_text(args.pptx)
    missing=[s for s in expected_strings(spec) if s not in text]
    assert not missing, f"renderer silently lost {len(missing)} values: {missing[:20]}"
    assert "PRESERVE_BODY_TOKEN_001" in text and "PRESERVE_BODY_TOKEN_120" in text
    assert slide_count>len(spec.get("slides",[]))+1, "dense fixture should create continuation slides"
    print(f"PASS no-silent-loss: {slide_count} PPT slides preserve all typed fixture content")
    return 0

if __name__=="__main__":
    raise SystemExit(main())
