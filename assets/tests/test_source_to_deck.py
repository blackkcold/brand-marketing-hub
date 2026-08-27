#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""End-to-end v4.1 source ingestion and coverage contract test."""
from __future__ import annotations
import json, subprocess, sys, tempfile
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from pptx import Presentation
from pypdf import PdfWriter
from PIL import Image as PILImage

ROOT=Path(__file__).resolve().parents[2]
INGEST=ROOT/"assets/ingest_sources.py"
COVER=ROOT/"assets/init_coverage.py"
VALID=ROOT/"assets/validate_v4_manifests.py"

def main():
    with tempfile.TemporaryDirectory(prefix="bmh-v41-") as d:
        d=Path(d)
        img=d/"source.png"; PILImage.new("RGB",(120,80),"white").save(img)
        doc=Document(); doc.add_heading("关键结论",1); doc.add_paragraph("必须保留的业务事实。"); doc.add_table(rows=2,cols=2); doc.add_picture(str(img)); doc.save(d/"a.docx")
        wb=Workbook(); ws=wb.active; ws.title="Data"; ws["A1"]="Metric"; ws["B1"]="Value"; ws["A2"]="Revenue"; ws["B2"]="=1+2"; ws.add_image(XLImage(str(img)),"D2"); wb.save(d/"b.xlsx")
        prs=Presentation(); s=prs.slides.add_slide(prs.slide_layouts[6]); box=s.shapes.add_textbox(0,0,1000000,1000000); box.text="旧汇报重点"; s.shapes.add_picture(str(img),1000000,1000000,width=1200000); prs.save(d/"c.pptx")
        w=PdfWriter(); w.add_blank_page(width=612,height=792); w.write(d/"d.pdf")
        (d/"e.csv").write_text("name,value\nA,1\n",encoding="utf-8")
        bundle=d/"bundle"
        p=subprocess.run([sys.executable,str(INGEST),str(d/"a.docx"),str(d/"b.xlsx"),str(d/"c.pptx"),str(d/"d.pdf"),str(d/"e.csv"),"--out-dir",str(bundle)],capture_output=True,text=True)
        assert p.returncode==0,p.stdout+p.stderr
        sources=json.loads((bundle/"source_inventory.json").read_text(encoding="utf-8"))
        content=json.loads((bundle/"content_units.json").read_text(encoding="utf-8"))
        assert len(sources["sources"])==5
        assert {s["kind"] for s in sources["sources"]}=={"docx","xlsx","pptx","pdf","csv"}
        ppt_source=next(s for s in sources["sources"] if s["kind"]=="pptx")
        assert ppt_source["role"]=="content-source", ppt_source
        assert any(u["unit_type"]=="cell-range" and u.get("formulae") for u in content["units"])
        assert any(u["unit_type"]=="slide-text" for u in content["units"])
        assert any(u["unit_type"]=="pdf-page" for u in content["units"])
        images=[u for u in content["units"] if u["unit_type"]=="image"]
        assert len(images)>=3, images
        for u in images:
            extracted=bundle/u["content"]["extracted_path"]
            assert extracted.exists() and extracted.stat().st_size>0
        assert all(not Path(s["path"]).is_absolute() for s in sources["sources"] if s.get("path"))

        styled=d/"styled"
        p=subprocess.run([sys.executable,str(INGEST),str(d/"c.pptx"),"--out-dir",str(styled),
                          "--role","c.pptx=style-reference"],capture_output=True,text=True)
        assert p.returncode==0,p.stdout+p.stderr
        styled_source=json.loads((styled/"source_inventory.json").read_text(encoding="utf-8"))["sources"][0]
        assert styled_source["role"]=="style-reference",styled_source

        first=content["units"][0]["unit_id"]
        deck={"version":"4.1","deck":{"title":"QA","purpose":"test","story_archetype":"custom","visual_template":"vivo-house",
              "template_manifest":"brand/vivo/template-manifest.json","runtime_preference":["pptxgenjs"],
              "source_of_truth":"deck_spec.json","confidentiality":"internal"},
              "slides":[{"slide_id":"S01","intent":"test","takeaway":"test","source_unit_ids":[first],
                         "content":[{"block_id":"BLK-001","type":"body_text","text":"test","unit_ids":[first]}],
                         "layout_candidates":["body"]}]}
        (d/"deck_spec.json").write_text(json.dumps(deck,ensure_ascii=False,indent=2),encoding="utf-8")
        cov=d/"coverage.json"
        p=subprocess.run([sys.executable,str(COVER),"--content",str(bundle/"content_units.json"),"--deck",str(d/"deck_spec.json"),"--out",str(cov)],capture_output=True,text=True)
        assert p.returncode==1, "coverage must fail because most source units are unmapped"
        coverage=json.loads(cov.read_text(encoding="utf-8"))
        assert coverage["summary"]["critical_missing"]>0

        # A fully mapped deck must pass coverage generation and cross-manifest validation.
        all_units=[u["unit_id"] for u in content["units"]]
        full_deck={"version":"4.1","deck":{"title":"QA Full","purpose":"test complete mapping","story_archetype":"custom",
                   "visual_template":"vivo-house","template_manifest":"brand/vivo/template-manifest.json",
                   "runtime_preference":["pptxgenjs"],"source_of_truth":"deck_spec.json","confidentiality":"internal"},
                   "slides":[{"slide_id":"S01","intent":"cover source","takeaway":"all source units are accounted for",
                              "source_unit_ids":all_units,
                              "content":[{"block_id":"BLK-001","type":"body_text","text":"all mapped","unit_ids":all_units}],
                              "layout_candidates":["body"]}]}
        full_deck_path=d/"deck_full.json"
        full_deck_path.write_text(json.dumps(full_deck,ensure_ascii=False,indent=2),encoding="utf-8")
        full_cov=d/"coverage_full.json"
        p=subprocess.run([sys.executable,str(COVER),"--content",str(bundle/"content_units.json"),
                          "--deck",str(full_deck_path),"--out",str(full_cov)],capture_output=True,text=True)
        assert p.returncode==0,p.stdout+p.stderr
        p=subprocess.run([sys.executable,str(VALID),
                          "--sources",str(bundle/"source_inventory.json"),
                          "--content",str(bundle/"content_units.json"),
                          "--coverage",str(full_cov),
                          "--deck",str(full_deck_path),
                          "--template",str(ROOT/"brand/vivo/template-manifest.json")],
                         capture_output=True,text=True)
        assert p.returncode==0,p.stdout+p.stderr
    print("PASS v4.1 source ingestion + fail/pass coverage + cross-manifest validation")
    return 0

if __name__=="__main__": raise SystemExit(main())
