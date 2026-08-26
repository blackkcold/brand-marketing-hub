#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
md2pptx.py — 提案骨架 markdown → pptx 骨架初稿转换器

规范：沿用 brand-marketing-hub 的通用 Markdown deck 方言；vivo 场景请使用 md2pptx_vivo.py。
  # 文档标题            → 封面页主标题
  <!-- deck: 副标题 -->  → 封面副标题
  ## P{n}｜页面标题      → 一页slide
  - 要点                → 一级bullet
    - 要点              → 二级bullet（缩进2空格）
  > 填写提示            → 演讲者备注（不出现在slide上）

用法：
  python3 md2pptx.py 输入.md 输出.pptx [--theme dark]
依赖：python3 -m pip install python-pptx --user
"""
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("缺少依赖：请先运行 python3 -m pip install python-pptx --user")

SLIDE_RE = re.compile(r'^##\s+P(\d+)[｜|]\s*(.+)$')
L1_RE = re.compile(r'^-\s+(.+)$')
L2_RE = re.compile(r'^\s{2,}-\s+(.+)$')
NOTE_RE = re.compile(r'^>\s*(.*)$')
DECK_RE = re.compile(r'<!--\s*deck:\s*(.+?)\s*-->')

# 主题
THEMES = {
    'light': {'bg': 'FFFFFF', 'title': '1F3864', 'text': '333333', 'accent': '4157FF'},
    'dark': {'bg': '111318', 'title': 'FFFFFF', 'text': 'D9DDE3', 'accent': '6C8CFF'},
}


def parse_md(path):
    """解析md为结构化页面列表"""
    pages = []
    deck_title, deck_sub = None, None
    current = None

    with open(path, encoding='utf-8') as f:
        for raw in f:
            line = raw.rstrip('\n')

            m = DECK_RE.search(line)
            if m and current is None:
                deck_sub = m.group(1)
                continue

            m = SLIDE_RE.match(line)
            if m:
                current = {'no': int(m.group(1)), 'title': m.group(2).strip(),
                           'bullets': [], 'notes': []}
                pages.append(current)
                continue

            if current is None:
                # 封面区（第一个 ## 之前）
                if line.startswith('# ') and not line.startswith('##'):
                    deck_title = line[2:].strip()
                continue

            m = NOTE_RE.match(line)
            if m and m.group(1) != '':
                current['notes'].append(m.group(1))
                continue

            m = L2_RE.match(line)
            if m:
                current['bullets'].append((1, m.group(1).strip()))
                continue

            m = L1_RE.match(line)
            if m:
                current['bullets'].append((0, m.group(1).strip()))
                continue

    return deck_title, deck_sub, pages


def add_textbox(slide, text, left, top, width, height, size, color,
                bold=False, align_center=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = 'PingFang SC'
    run.font.color.rgb = RGBColor.from_string(color)
    return tb


def build_pptx(deck_title, deck_sub, pages, out_path, theme):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 封面页：若模板首个P页标题为"封面"，则将其渲染为封面版式；否则自动生成封面
    body_pages = list(pages)
    has_cover_page = bool(pages) and pages[0]['title'].strip() == '封面'
    if has_cover_page:
        cover_pg = pages[0]
        body_pages = pages[1:]
        cover = prs.slides.add_slide(blank)
        add_textbox(cover, deck_title or '未命名方案', 1.0, 1.8, 11.3, 1.2, 40,
                    theme['title'], bold=True, align_center=True)
        if deck_sub:
            add_textbox(cover, deck_sub, 1.0, 3.1, 11.3, 0.7, 20,
                        theme['accent'], align_center=True)
        # 封面页的要点居中排在下方（提报方/日期/密级等）
        for i, (_, text) in enumerate(cover_pg['bullets']):
            add_textbox(cover, text, 1.0, 4.2 + i * 0.55, 11.3, 0.5, 15,
                        theme['text'], align_center=True)
        if cover_pg['notes']:
            cover.notes_slide.notes_text_frame.text = '填写提示：' + ' '.join(cover_pg['notes'])
    else:
        cover = prs.slides.add_slide(blank)
        add_textbox(cover, deck_title or '未命名方案', 1.0, 2.6, 11.3, 1.2, 40,
                    theme['title'], bold=True, align_center=True)
        if deck_sub:
            add_textbox(cover, deck_sub, 1.0, 4.0, 11.3, 0.7, 20,
                        theme['accent'], align_center=True)

    # 内容页
    for pg in body_pages:
        slide = prs.slides.add_slide(blank)
        # 标题
        add_textbox(slide, pg['title'], 0.7, 0.4, 12.0, 0.9, 28,
                    theme['title'], bold=True)
        # 标题下强调线
        line = slide.shapes.add_shape(1, Inches(0.75), Inches(1.35),
                                      Inches(1.6), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(theme['accent'])
        line.line.fill.background()

        # 要点
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.6), Inches(5.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, (level, text) in enumerate(pg['bullets']):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = level
            run = p.add_run()
            run.text = ('• ' if level == 0 else '– ') + text
            run.font.size = Pt(18 if level == 0 else 15)
            run.font.name = 'PingFang SC'
            run.font.color.rgb = RGBColor.from_string(
                theme['text'] if level == 0 else theme['accent'])
            p.space_after = Pt(8 if level == 0 else 4)

        # 演讲者备注
        if pg['notes']:
            notes = slide.notes_slide.notes_text_frame
            notes.text = '填写提示：' + ' '.join(pg['notes'])

    prs.save(out_path)
    return len(prs.slides._sldIdLst)  # 总页数


def main():
    argv = sys.argv[1:]
    theme_name = 'light'
    args = []
    i = 0
    while i < len(argv):
        if argv[i] == '--theme':
            i += 1
            if i < len(argv):
                theme_name = argv[i] if argv[i] in THEMES else 'light'
        elif argv[i].startswith('--'):
            if argv[i] == '--dark':
                theme_name = 'dark'
        else:
            args.append(argv[i])
        i += 1
    if len(args) != 2:
        sys.exit(__doc__)
    theme = THEMES[theme_name]

    src, dst = args
    title, sub, pages = parse_md(src)

    # 转换前结构校验
    nos = [p['no'] for p in pages]
    if nos != list(range(1, len(pages) + 1)):
        print(f'[警告] 页码不连续或未从P1开始: {nos}')
    for p in pages:
        l1 = sum(1 for lv, _ in p['bullets'] if lv == 0)
        if l1 > 6:
            print(f'[警告] P{p["no"]}「{p["title"]}」一级要点 {l1} 条，超过6条可能溢出')

    total = build_pptx(title, sub, pages, dst, theme)
    print(f'[完成] {dst}')
    print(f'  内容页: {len(pages)} (P1–P{len(pages)})，加封面共 {total} 页')
    print(f'  主题: {theme_name} | 中文字体: PingFang SC')


if __name__ == '__main__':
    main()
