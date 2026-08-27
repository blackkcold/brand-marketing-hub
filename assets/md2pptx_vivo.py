#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""vivo 2026 Markdown -> PPTX renderer.

The renderer is intentionally content-led: the Markdown layout directive selects
an archetype instead of forcing every page into a blue title bar + bullets page.
Old syntax remains supported for backwards compatibility.
"""
from __future__ import annotations

import glob
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:
    Image = None

SW, SH = 13.333, 7.5
BLUE, BRIGHT, NAVY = '1E46E6', '3458F6', '06175E'
LIGHT, SURFACE, INK = 'D1EBFE', 'EAF2FF', '111111'
GRAY, MUTED, GOLD, RED, WHITE = '565656', '374151', 'F6C84C', 'E6001E', 'FFFFFF'
GRAD_A, GRAD_B = '3458F6', '1E46E6'
TABLE_ALT = 'F2F7FF'
LOGO = Path(__file__).resolve().parent / 'vivo-deck' / 'vivo_wordmark_white.png'
HOUSE_TEMPLATE = Path(os.environ.get(
    'VIVO_PPT_TEMPLATE',
    '/Users/11169285/Library/CloudStorage/OneDrive-个人/桌面/vivo/产品 - 品牌营销/vivo PPT模版_20230727-20250325.pptx',
))
OVERFLOW_WARNINGS: list[str] = []

PROFILE_PALETTES = {
    'vivo-house': (BLUE, BRIGHT, NAVY, LIGHT),
    'campaign': (BLUE, BRIGHT, NAVY, LIGHT),
    'celebrity': (BLUE, BRIGHT, NAVY, LIGHT),
    'ip-collab': (BLUE, BRIGHT, NAVY, LIGHT),
    'first': ('4160FF', '3756FE', NAVY, LIGHT),
    'disney': ('4661F4', '80A0D0', NAVY, 'E0C0E0'),
    'iqoo': ('C00000', 'F0B419', '111111', 'FFF4D6'),
    'dark-trend': ('111111', '5B4BFF', 'FFFFFF', '242424'),
}

FONT_CHAIN = ['微软雅黑']


def C(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace('#', ''))


def font_available(name: str) -> bool:
    if name in ('vivoSans', 'VIVOTYPECN'):
        patterns = [f'/Library/Fonts/*{name}*', f'/System/Library/Fonts/**/*{name}*',
                    os.path.expanduser(f'~/Library/Fonts/*{name}*')]
        return any(glob.glob(p, recursive=True) for p in patterns)
    if name == 'PingFang SC':
        return os.path.exists('/System/Library/Fonts/PingFang.ttc')
    if name == '微软雅黑':
        return os.path.exists('/Library/Fonts/msyh.ttf') or sys.platform.startswith('win')
    return True


# The output contract is strict even when the local QA renderer substitutes the font.
FONT = '微软雅黑'


def set_font(run, name: str = FONT):
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        node = rpr.find(qn(tag))
        if node is None:
            node = rpr.makeelement(qn(tag), {})
            rpr.append(node)
        node.set('typeface', name)


def replace_text(shape, text, size, color=WHITE, bold=False):
    """Replace a template placeholder while keeping it as an editable text shape."""
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.line_spacing = 1.15
    r = p.add_run(); r.text = str(text); r.font.size = Pt(size)
    r.font.color.rgb = C(color); r.font.bold = bold; set_font(r)


def scale_template_slide(slide, sx, sy):
    layout_ph = {p.placeholder_format.idx: p for p in slide.slide_layout.placeholders}
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp_pr = shape._element.find(qn('p:spPr'))
            xfrm = sp_pr.find(qn('a:xfrm')) if sp_pr is not None else None
            if xfrm is None:
                src = layout_ph.get(shape.placeholder_format.idx)
                if src is not None and src._element.find(qn('p:spPr')).find(qn('a:xfrm')) is not None:
                    from copy import deepcopy
                    shape._element.find(qn('p:spPr')).append(deepcopy(src._element.find(qn('p:spPr')).find(qn('a:xfrm'))))
        shape.left = int(shape.left * sx); shape.top = int(shape.top * sy)
        shape.width = int(shape.width * sx); shape.height = int(shape.height * sy)


def scale_template_layout(layout, sx, sy):
    """Scale background pictures and shapes on a layout to fill the new canvas."""
    sp_tree = layout._element.find(qn('p:cSld') + '/' + qn('p:spTree'))
    if sp_tree is None:
        return
    for xfrm in sp_tree.iter(qn('a:xfrm')):
        off = xfrm.find(qn('a:off')); ext = xfrm.find(qn('a:ext'))
        if off is not None and off.get('x') is not None:
            off.set('x', str(int(round(float(off.get('x')) * sx))))
        if off is not None and off.get('y') is not None:
            off.set('y', str(int(round(float(off.get('y')) * sy))))
        if ext is not None and ext.get('cx') is not None:
            ext.set('cx', str(int(round(float(ext.get('cx')) * sx))))
        if ext is not None and ext.get('cy') is not None:
            ext.set('cy', str(int(round(float(ext.get('cy')) * sy))))


def prepare_house_template(deck):
    """Load the requested vivo cover/end pages and keep them native/editable."""
    prs = Presentation(str(HOUSE_TEMPLATE))
    source_w, source_h = prs.slide_width / 914400, prs.slide_height / 914400
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    sx, sy = SW / source_w, SH / source_h
    cover_slide, end_slide = prs.slides[0], prs.slides[-1]
    # Reserve a high slide part name so python-pptx can append generated slides
    # without creating duplicate ZIP members, while preserving the native end page.
    end_slide.part._partname = PackURI('/ppt/slides/slide99.xml')
    scale_template_slide(cover_slide, sx, sy); scale_template_slide(end_slide, sx, sy)
    scale_template_layout(cover_slide.slide_layout, sx, sy)
    scale_template_layout(end_slide.slide_layout, sx, sy)
    text_shapes = [s for s in cover_slide.shapes if getattr(s, 'has_text_frame', False)]
    placeholders = [s for s in cover_slide.placeholders if getattr(s, 'has_text_frame', False)]
    if len(placeholders) >= 2:
        replace_text(placeholders[0], deck.get('title') or '未命名 deck', 40, WHITE, True)
        replace_text(placeholders[1], deck.get('sub') or '', 20, 'D1EBFE')
    if text_shapes:
        replace_text(text_shapes[0], deck.get('meta') or '', 9, WHITE)
    if len(text_shapes) > 1:
        replace_text(text_shapes[1], '保密等级', 9, 'F6C84C')
    end_placeholders = [s for s in end_slide.placeholders if getattr(s, 'has_text_frame', False)]
    if end_placeholders:
        replace_text(end_placeholders[0], 'THANK YOU\n谢谢', 40, WHITE, True)
    return prs, end_slide


def move_slide_to_end(prs, slide):
    for sld_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(sld_id.rId) is slide.part:
            prs.slides._sldIdLst.remove(sld_id)
            prs.slides._sldIdLst.append(sld_id)
            return


def textbox(slide, x, y, w, h, text, size=12, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04,
            font=FONT, italic=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin); tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = 1.15
    run = p.add_run(); run.text = str(text); run.font.size = Pt(size)
    run.font.color.rgb = C(color); run.font.bold = bold; run.font.italic = italic; set_font(run, font)
    return shape


def fit_text_size(text, width, height, preferred, minimum=10, line_factor=1.18):
    """Return the largest readable size that fits the requested text box."""
    size = preferred
    while size > minimum and est_lines(text, size, width) * size * line_factor / 72 > height:
        size -= 0.5
    return max(size, minimum)


def rect(slide, x, y, w, h, fill, radius=False, line=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = C(fill)
    if transparency:
        shape.fill.transparency = transparency / 100
    if line:
        shape.line.color.rgb = C(line)
    else:
        shape.line.fill.background()
    return shape


def cjk_width(text, size):
    return sum(size if ord(ch) > 0x2E00 else size * 0.52 for ch in text) / 72


def est_lines(text, size, width):
    total = 0
    for para in str(text).split('\n'):
        current = 0
        lines = 1
        for ch in para:
            current += size if ord(ch) > 0x2E00 else size * 0.52
            if current / 72 > max(width, 0.2):
                lines += 1; current = size if ord(ch) > 0x2E00 else size * 0.52
        total += lines
    return max(total, 1)


def warn(page, label, height, box_height):
    if height > box_height + 0.01:
        OVERFLOW_WARNINGS.append(f'P{page["no"]}「{page["title"]}」{label}: {height:.2f}in > {box_height:.2f}in')


SLIDE_RE = re.compile(r'^##\s+P(\d+)[｜|]\s*(.+)$')


def parse_md(path):
    deck = {'title': None, 'sub': None, 'meta': None, 'profile': 'vivo-house', 'pages': [], 'source': path}
    cur = None
    for raw in Path(path).read_text(encoding='utf-8').splitlines():
        line = raw.rstrip()
        comment = re.search(r'<!--\s*deck:\s*(.+?)\s*-->', line)
        if comment and cur is None:
            body = comment.group(1).strip()
            if body.startswith('meta:'): deck['meta'] = body[5:].strip()
            else: deck['sub'] = body
            continue
        if line.startswith('# ') and cur is None:
            deck['title'] = line[2:].strip(); continue
        match = SLIDE_RE.match(line)
        if match:
            cur = {'no': int(match.group(1)), 'title': match.group(2).strip(), 'layout': None,
                   'arg': None, 'sub': None, 'source': None, 'profile': None, 'img': [],
                    'bullets': [], 'cards': [], 'table': [], 'notes': [], 'stats': [], 'chart': None}
            deck['pages'].append(cur); continue
        if cur is None: continue
        if line.startswith('@layout'):
            parts = line.split(); cur['layout'] = parts[1] if len(parts) > 1 else 'bullets'; cur['arg'] = ' '.join(parts[2:])
        elif line.startswith('@profile'):
            cur['profile'] = line.split(maxsplit=1)[1].strip(); deck['profile'] = cur['profile']
        elif line.startswith('@sub'): cur['sub'] = line[4:].strip()
        elif line.startswith('@source'): cur['source'] = line[7:].strip()
        elif line.startswith('@img'):
            rest = line[4:].strip(); p, _, cap = rest.partition('|'); cur['img'].append((p.strip(), cap.strip()))
        elif line.startswith('@stat'):
            rest = line[5:].strip(); value, _, label = rest.partition('|'); cur['stats'].append((value.strip(), label.strip()))
        elif line.startswith('@chart'):
            chart_type = line.split(maxsplit=1)[1].strip().lower() if len(line.split(maxsplit=1)) > 1 else 'bar'
            aliases = {'column': 'bar', 'doughnut': 'doughnut', 'donut': 'doughnut', '环图': 'doughnut', '柱状图': 'bar', '折线图': 'line'}
            cur['chart'] = aliases.get(chart_type, chart_type)
        elif line.startswith('### '): cur['cards'].append({'title': line[4:].strip(), 'items': []})
        elif line.startswith('> '): cur['notes'].append(line[2:].strip())
        elif line.strip().startswith('|'):
            cells = [c.strip() for c in line.strip().strip('|').split('|')]
            if not all(set(c) <= set('-: ') for c in cells): cur['table'].append(cells)
        elif re.match(r'^\s{2,}-\s+', line):
            item = re.sub(r'^\s{2,}-\s+', '', line)
            if cur['cards']: cur['cards'][-1]['items'].append(item)
            else: cur['bullets'].append((1, item))
        elif re.match(r'^-\s+', line):
            item = re.sub(r'^-\s+', '', line)
            if cur['cards']: cur['cards'][-1]['items'].append(item)
            else: cur['bullets'].append((0, item))
    return deck


def gradient(slide, profile='vivo-house'):
    primary, bright, _, _ = PROFILE_PALETTES.get(profile, PROFILE_PALETTES['vivo-house'])
    bg = rect(slide, 0, 0, SW, SH, primary)
    try:
        bg.fill.gradient(); stops = bg.fill.gradient_stops
        stops[0].color.rgb = C(bright); stops[1].color.rgb = C(primary)
    except Exception:
        pass
    for x, w, transparency in ((10.1, 1.35, 90), (11.25, .8, 93), (9.35, .5, 95)):
        band = rect(slide, x, -2.5, w, 14, WHITE, transparency=transparency); band.rotation = 24


def add_logo(slide):
    if LOGO.exists(): slide.shapes.add_picture(str(LOGO), Inches(11.65), Inches(.42), height=Inches(.42))


def chrome(slide, page, total, source=None):
    textbox(slide, .58, 7.04, 7.5, .25, source or '内部汇报资料 · 请勿外传', 8.5, MUTED)
    textbox(slide, 12.0, 7.04, .72, .25, str(page['no']), 8.5, MUTED, align=PP_ALIGN.RIGHT)


def cover(prs, deck, pg=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); gradient(slide, 'vivo-house'); add_logo(slide)
    cover_title = deck['title'] or '未命名 deck'
    title_h = 1.65 if est_lines(cover_title, 40, 8.45) <= 2 else 2.1
    title_size = fit_text_size(cover_title, 8.45, title_h, 40, 30, 1.15)
    textbox(slide, .9, 2.05, 8.45, title_h, cover_title, title_size, WHITE, True, margin=0)
    if deck['sub']: textbox(slide, .92, 4.05 if title_h <= 1.65 else 4.35, 7.25, .6, deck['sub'], 20, LIGHT, margin=0)
    if deck['meta']: textbox(slide, .92, 6.48, 10, .35, deck['meta'], 9, WHITE, margin=0)
    return slide


def end(slide, profile='vivo-house'):
    gradient(slide, 'vivo-house'); add_logo(slide)
    textbox(slide, .9, 2.85, 10, .8, 'THANK YOU', 44, WHITE, True, margin=0)
    textbox(slide, .92, 3.85, 10, .5, '谢谢', 20, LIGHT, margin=0)


def title(slide, page):
    textbox(slide, .62, .48, 8.65, .58, page['title'], 22, NAVY, True, margin=0)
    if page['sub']: textbox(slide, 9.45, .54, 3.25, .32, page['sub'], 11, GRAY, True, align=PP_ALIGN.RIGHT, margin=0)


def bullets(slide, page, x=.7, y=1.45, w=11.9, h=5.1, size=15, with_title=True):
    if with_title:
        title(slide, page)
    content = '\n'.join(('• ' if level == 0 else '   – ') + text for level, text in page['bullets'])
    size = fit_text_size(content, w - .08, h, size, 12, 1.18)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = Inches(.04); tf.margin_right = Inches(.04)
    for i, (level, text) in enumerate(page['bullets']):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.line_spacing = 1.18; p.space_after = Pt(9 if level == 0 else 4)
        r = p.add_run(); r.text = ('• ' if level == 0 else '   – ') + text; r.font.size = Pt(size - level)
        r.font.color.rgb = C(INK if level == 0 else GRAY); r.font.bold = level == 0; set_font(r)
    warn(page, '正文', sum(est_lines(t, size-level, w) * (size-level) * 1.18 / 72 for level, t in page['bullets']), h)


def toc(slide, page):
    textbox(slide, .7, .62, 4, .6, '目录', 28, NAVY, True, margin=0)
    y = 1.55
    for i, (_, text) in enumerate(page['bullets'], 1):
        parts = text.split(maxsplit=1); num = parts[0] if parts and parts[0].isdigit() else f'{i:02d}'
        label = parts[1] if len(parts) > 1 else text
        textbox(slide, .9, y, .65, .42, num.zfill(2), 23, BLUE, True, margin=0)
        textbox(slide, 1.75, y, 9.8, .42, label, 21, INK, margin=0)
        rect(slide, 1.75, y + .55, 9.7, .018, BLUE)
        y += .88


def part(slide, page):
    raw = page['title']; match = re.search(r'(\d+)', raw); num = match.group(1) if match else ''
    label = re.sub(r'^Part\s*\d*\s*', '', raw, flags=re.I).strip()
    textbox(slide, .9, 1.2, 2.4, .5, f'Part {num}', 25, INK, True, margin=0)
    textbox(slide, .9, 2.0, 3.6, 1.0, num or '—', 80, BLUE, True, margin=0)
    textbox(slide, 4.15, 2.42, 8.3, .9, label, 40, BLUE, True, margin=0)
    if page['sub']: textbox(slide, 4.18, 3.48, 8, .45, page['sub'], 18, GRAY, margin=0)


def cards(slide, page, ncol=3):
    title(slide, page); data = page['cards'] or [{'title': '', 'items': [t for _, t in page['bullets']]}]
    ncol = max(1, min(ncol, 4, len(data))); rows = (len(data) + ncol - 1) // ncol
    gap, left, top = .28, .65, 1.42; cw = (SW - 1.3 - gap * (ncol - 1)) / ncol
    rh = (5.15 - gap * (rows - 1)) / rows
    for i, card in enumerate(data):
        x = left + (i % ncol) * (cw + gap); y = top + (i // ncol) * (rh + gap)
        rect(slide, x, y, cw, rh, SURFACE, radius=True)
        textbox(slide, x + .22, y + .18, cw - .44, .42, card['title'] or '要点', 15, BLUE, True, margin=0)
        body = textbox(slide, x + .22, y + .75, cw - .44, rh - .95, '\n'.join('• ' + t for t in card['items']), 11.5, INK, margin=0)
        body.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        warn(page, f'卡片{i + 1}', sum(est_lines(t, 11.5, cw - .44) * 11.5 * 1.15 / 72 for t in card['items']), rh - .95)


def stats(slide, page):
    title(slide, page); data = page['stats'] or [(c['title'], c['items'][0] if c['items'] else '') for c in page['cards']]
    if not data: data = [('关键数据', t) for _, t in page['bullets']]
    n = max(1, min(4, len(data))); gap = .3; cw = (12.0 - gap * (n - 1)) / n
    for i, (value, label) in enumerate(data[:4]):
        card_h = 3.0 if len(data) <= 2 else 3.75
        y = 1.8 if len(data) <= 2 else 1.65
        x = .65 + i * (cw + gap); rect(slide, x, y, cw, card_h, BLUE if i == 0 else LIGHT, radius=True)
        color = WHITE if i == 0 else NAVY
        value_size = fit_text_size(value, cw - .44, .95, 40, 24, 1.05)
        textbox(slide, x + .22, y + .52, cw - .44, .95, value, value_size, color, True, margin=0)
        textbox(slide, x + .22, y + 1.72, cw - .44, card_h - 1.95, label, 14, color, margin=0)


def framework(slide, page):
    title(slide, page); data = page['cards'][:4]
    if not data: data = [{'title': 'RTB', 'items': [t for _, t in page['bullets']]}]
    n = min(4, len(data)); gap = .25; cw = (12 - gap * (n - 1)) / n
    rh = 4.6
    for i, card in enumerate(data):
        x = .65 + i * (cw + gap)
        rect(slide, x, 1.55, cw, rh, LIGHT, radius=True)
        textbox(slide, x + .2, 1.9, cw - .4, .5, card['title'].upper(), 16, BLUE, True, margin=0)
        body = textbox(slide, x + .2, 2.55, cw - .4, rh - 1.2, '\n'.join('• ' + t for t in card['items']), 12, INK, margin=0)
        body.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def table(slide, page):
    title(slide, page); rows = page['table'] or [['项目', '内容'], *[[t, '待补'] for _, t in page['bullets']]]
    cols = max(len(r) for r in rows); rows = [r + [''] * (cols - len(r)) for r in rows]; rows = rows[:10]; cols = min(cols, 6)
    x, y, w, h = .58, 1.45, 12.18, min(5.25, .46 * len(rows) + .2)
    shape = slide.shapes.add_table(len(rows), cols, Inches(x), Inches(y), Inches(w), Inches(h)); tbl = shape.table
    widths = [w / cols] * cols
    if cols >= 4: widths[0] = w * .16; widths[1] = w * .22; rem = w - widths[0] - widths[1]; widths[2:] = [rem / (cols - 2)] * (cols - 2)
    for i, width in enumerate(widths): tbl.columns[i].width = Inches(width)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(.46 if ri == 0 else .53)
        for ci in range(cols):
            cell = tbl.cell(ri, ci); cell.text = row[ci]; cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = C(BLUE if ri == 0 else (TABLE_ALT if ri % 2 == 0 else WHITE))
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(10.5); r.font.bold = ri == 0 or ci == 0; r.font.color.rgb = C(WHITE if ri == 0 else (NAVY if ci == 0 else INK)); set_font(r)


def comparison(slide, page):
    title(slide, page); data = page['cards'][:2]
    if len(data) < 2: return cards(slide, page, 2)
    rh = 4.6
    for i, card in enumerate(data):
        x = .7 + i * 6.15; rect(slide, x, 1.55, 5.75, rh, LIGHT if i else SURFACE, radius=True)
        textbox(slide, x + .3, 1.95, 5.15, .62, card['title'], 19, BLUE if i else NAVY, True, margin=0)
        body = textbox(slide, x + .3, 2.75, 5.1, rh - 1.35, '\n'.join('• ' + t for t in card['items']), 14, INK, margin=0)
        body.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def timeline(slide, page):
    title(slide, page); rows = page['table'] or [['阶段', '时间', '关键动作'], *[[t, '', '待补'] for _, t in page['bullets']]]
    rows = rows[:7]; y = 1.65
    for i, row in enumerate(rows[1:]):
        label = row[0] if row else ''; time = row[1] if len(row) > 1 else ''; action = ' · '.join(row[2:]) if len(row) > 2 else ''
        rect(slide, 1.2, y + .48, 10.8, .045, 'AAB9D5'); rect(slide, 1.2, y, 1.45, .42, BLUE, radius=True)
        textbox(slide, 1.35, y + .07, 1.15, .26, label, 11, WHITE, True, margin=0)
        rect(slide, 2.95, y + .39, .2, .2, RED if '截止' in time or 'deadline' in time.lower() else BLUE, radius=True)
        textbox(slide, 3.35, y - .02, 1.65, .3, time, 11, NAVY, True, margin=0)
        textbox(slide, 5.05, y - .02, 6.7, .38, action, 12, INK, margin=0)
        y += .72


def collage(slide, page):
    title(slide, page); imgs = page['img'][:6]
    if not imgs: return split(slide, page, 'left')
    if not any(Path(path).exists() for path, _ in imgs):
        return bullets(slide, page, y=1.55, h=4.85, size=13, with_title=False)
    n = len(imgs); cols = 3 if n > 2 else n; gap = .14; cw = (8.0 - gap * (cols - 1)) / cols; rh = 4.7 if n <= cols else 2.25
    for i, (path, cap) in enumerate(imgs):
        x = .65 + (i % cols) * (cw + gap); y = 1.45 + (i // cols) * (rh + gap)
        if Path(path).exists(): add_cover_crop(slide, path, x, y, cw, rh)
        else: rect(slide, x, y, cw, rh, LIGHT, line=BLUE); textbox(slide, x, y + rh / 2 - .2, cw, .3, 'MISSING IMAGE', 10, BLUE, True, align=PP_ALIGN.CENTER)
    textbox(slide, 9.05, 1.65, 3.5, 4.4, '\n'.join('• ' + t for _, t in page['bullets']), 14, INK, margin=0)


def chart(slide, page):
    title(slide, page)
    rows = page['table']
    if len(rows) < 2 or len(rows[0]) < 2:
        return bullets(slide, page, y=1.55, h=4.85, size=13, with_title=False)
    data = ChartData()
    data.categories = [row[0] for row in rows[1:] if row]
    for col in range(1, len(rows[0])):
        name = rows[0][col] or f'系列 {col}'
        values = []
        for row in rows[1:]:
            try: values.append(float(row[col]))
            except (ValueError, IndexError): values.append(0)
        data.add_series(name, values)
    kind = {'bar': XL_CHART_TYPE.COLUMN_CLUSTERED, 'line': XL_CHART_TYPE.LINE_MARKERS,
            'doughnut': XL_CHART_TYPE.DOUGHNUT}.get(page.get('chart'), XL_CHART_TYPE.COLUMN_CLUSTERED)
    graphic = slide.shapes.add_chart(kind, Inches(.85), Inches(1.55), Inches(11.7), Inches(4.75), data)
    ch = graphic.chart
    ch.has_legend = len(data) > 1 or page.get('chart') == 'doughnut'
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    ch.has_title = False
    try:
        ch.value_axis.has_major_gridlines = False
        ch.category_axis.tick_labels.font.name = FONT
        ch.value_axis.tick_labels.font.name = FONT
        ch.category_axis.tick_labels.font.size = Pt(10)
        ch.value_axis.tick_labels.font.size = Pt(10)
    except (AttributeError, ValueError):
        pass
    for i, series in enumerate(ch.series):
        series.format.fill.solid(); series.format.fill.fore_color.rgb = C((BLUE, BRIGHT, NAVY, GOLD)[i % 4])


def add_cover_crop(slide, path, x, y, w, h):
    if Image is None: return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    iw, ih = Image.open(path).size; src, dst = iw / ih, w / h
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if src > dst:
        crop = (1 - dst / src) / 2; pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - src / dst) / 2; pic.crop_top = crop; pic.crop_bottom = crop
    return pic


def add_contain(slide, path, x, y, w, h, background=None):
    """Fit the full product/person into a frame; case-study evidence must not crop the hero."""
    if background:
        rect(slide, x, y, w, h, background)
    if Image is None:
        return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    iw, ih = Image.open(path).size; scale = min(w / iw, h / ih)
    nw, nh = iw * scale, ih * scale
    return slide.shapes.add_picture(path, Inches(x + (w - nw) / 2), Inches(y + (h - nh) / 2),
                                    width=Inches(nw), height=Inches(nh))


def split(slide, page, side='left'):
    title(slide, page); img_w, img_h = 5.25, 4.85; ix = .65 if side != 'right' else 7.43
    if page['img']:
        path, cap = page['img'][0]
        if Path(path).exists():
            source_ratio = img_w / img_h
            if Image is not None:
                with Image.open(path) as image:
                    source_ratio = image.size[0] / image.size[1]
            frame_w, frame_h = img_w, img_h
            if source_ratio > img_w / img_h:
                frame_h = img_w / source_ratio
            elif source_ratio < img_w / img_h:
                frame_w = img_h * source_ratio
            frame_x, frame_y = ix + (img_w - frame_w) / 2, 1.5 + (img_h - frame_h) / 2
            add_contain(slide, path, frame_x, frame_y, frame_w, frame_h, WHITE)
        else: rect(slide, ix, 1.5, img_w, img_h, LIGHT, line=BLUE); textbox(slide, ix, 3.7, img_w, .3, 'MISSING IMAGE', 11, BLUE, True, align=PP_ALIGN.CENTER)
        if cap: textbox(slide, ix, 6.45, img_w, .25, cap, 8.5, MUTED, align=PP_ALIGN.CENTER, margin=0)
    bx, bw = (.65, 6.05) if side == 'right' else (6.45, 6.1)
    bullets(slide, page, bx, 1.55, bw, 4.85, 13, with_title=False)


def expand_overflow_pages(pages):
    """Split dense bullet pages only after keeping the readable minimum size."""
    expanded = []
    for page in pages:
        layout = page.get('layout') or 'bullets'
        if layout not in ('bullets', 'purpose') or len(page.get('bullets', [])) < 2:
            expanded.append(page)
            continue
        chunks, current, height = [], [], 0.0
        for level, text in page['bullets']:
            item_height = est_lines(text, 15 - level, 11.75) * (15 - level) * 1.18 / 72 + (0.13 if level == 0 else 0.06)
            if current and height + item_height > 4.9:
                chunks.append(current); current = []; height = 0.0
            current.append((level, text)); height += item_height
        if current: chunks.append(current)
        if len(chunks) == 1:
            expanded.append(page)
            continue
        for index, chunk in enumerate(chunks, 1):
            split_page = dict(page)
            split_page['bullets'] = chunk
            split_page['title'] = f"{page['title']}（续 {index}/{len(chunks)}）" if index > 1 else page['title']
            split_page['notes'] = page.get('notes', []) if index == 1 else []
            expanded.append(split_page)
    return expanded


def build(deck, out):
    global OVERFLOW_WARNINGS; OVERFLOW_WARNINGS = []
    if HOUSE_TEMPLATE.exists():
        prs, end_slide = prepare_house_template(deck)
    else:
        prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH); end_slide = None
    blank = prs.slide_layouts[6]
    pages = list(deck['pages']); cover_pg = pages.pop(0) if pages and pages[0]['title'] == '封面' else None
    if end_slide is not None and pages and (pages[-1].get('layout') == 'end' or '谢谢' in pages[-1].get('title', '')):
        pages.pop()
    pages = expand_overflow_pages(pages)
    for number, page in enumerate(pages, 2):
        page['no'] = number
    if end_slide is None:
        cover(prs, deck, cover_pg)
    else:
        # The source file already contains the native cover and end pages.
        # Remove its instructional middle slide before adding generated content.
        if len(prs.slides) > 2:
            middle_id = prs.slides._sldIdLst[1]
            prs.slides._sldIdLst.remove(middle_id)
    total = len(pages) + 1
    for page in pages:
        layout = page['layout'] or ('toc' if page['title'] == '目录' else ('end' if 'thank' in page['title'].lower() or '谢谢' in page['title'] else 'bullets'))
        slide = prs.slides.add_slide(blank)
        if layout == 'end':
            if end_slide is None: end(slide, page.get('profile') or deck['profile'])
        elif layout == 'toc': toc(slide, page)
        elif layout == 'part': part(slide, page)
        elif layout == 'stats': stats(slide, page)
        elif layout in ('framework', 'funnel'): framework(slide, page)
        elif layout == 'cards': cards(slide, page, int(page['arg']) if page['arg'].isdigit() else 3)
        elif layout == 'comparison': comparison(slide, page)
        elif layout == 'matrix': table(slide, page)
        elif layout == 'timeline': timeline(slide, page)
        elif layout in ('budget', 'table'): table(slide, page)
        elif layout in ('collage', 'case-study'): collage(slide, page)
        elif layout == 'chart' or page.get('chart'): chart(slide, page)
        elif layout == 'split': split(slide, page, page['arg'] if page['arg'] in ('left', 'right') else 'left')
        else: bullets(slide, page)
        if layout != 'end': chrome(slide, page, total, page.get('source'))
        if page['notes']: slide.notes_slide.notes_text_frame.text = '\n'.join(page['notes'])
    if end_slide is not None:
        move_slide_to_end(prs, end_slide)
    prs.save(out); return len(prs.slides)


def main():
    if len(sys.argv) != 3: raise SystemExit('用法: python3 md2pptx_vivo.py 输入.md 输出.pptx')
    deck = parse_md(sys.argv[1]); nums = [p['no'] for p in deck['pages']]
    if nums != list(range(1, len(nums) + 1)): print(f'[警告] 页码不连续: {nums}')
    for p in deck['pages']:
        if sum(1 for level, _ in p['bullets'] if level == 0) > 6: print(f'[警告] P{p["no"]} 一级要点超过6条')
    count = build(deck, sys.argv[2]); print(f'[完成] {sys.argv[2]} 共 {count} 页'); print(f'[字体] {FONT}')
    for item in OVERFLOW_WARNINGS: print('[溢出] ' + item)


if __name__ == '__main__': main()
